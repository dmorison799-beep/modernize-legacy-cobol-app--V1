#!/usr/bin/env python3
"""
Generate PowerPoint presentation for COBOL to Node.js migration.
Run: python presentation/generate_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
LIGHT_BLUE = RGBColor(0x3A, 0x7C, 0xA5)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE6, 0x5C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background shape
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = DARK_BLUE
    bg_shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "COBOL to Node.js Modernization"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Legacy Accounting System Migration"
    p2.font.size = Pt(24)
    p2.font.color.rgb = LIGHT_BLUE
    p2.alignment = PP_ALIGN.CENTER

    # Description
    txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(5), Inches(7), Inches(1.5))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "Converting a multi-module COBOL accounting system to a modern Node.js application with full test coverage, Docker deployment, and cloud-ready infrastructure."
    p3.font.size = Pt(14)
    p3.font.color.rgb = WHITE
    p3.alignment = PP_ALIGN.CENTER

    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = (
        "Welcome to the COBOL to Node.js modernization presentation. "
        "This project demonstrates converting a legacy COBOL accounting system "
        "into a modern, cloud-native Node.js application. "
        "We'll cover the legacy system, our approach, code mapping, testing, and deployment."
    )


def add_section_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = LIGHT_BLUE
    bg_shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(7), Inches(1))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullets, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content
    txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)
        if bullet.startswith("  "):
            p.level = 1
            p.font.size = Pt(14)

    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide


def add_code_slide(prs, title, left_title, left_code, right_title, right_code, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Left column header
    left_hdr = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.4), Inches(4.5), Inches(0.5)
    )
    left_hdr.fill.solid()
    left_hdr.fill.fore_color.rgb = ORANGE
    left_hdr.line.fill.background()
    left_hdr.text_frame.paragraphs[0].text = left_title
    left_hdr.text_frame.paragraphs[0].font.size = Pt(12)
    left_hdr.text_frame.paragraphs[0].font.bold = True
    left_hdr.text_frame.paragraphs[0].font.color.rgb = WHITE
    left_hdr.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Right column header
    right_hdr = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(1.4), Inches(4.5), Inches(0.5)
    )
    right_hdr.fill.solid()
    right_hdr.fill.fore_color.rgb = GREEN
    right_hdr.line.fill.background()
    right_hdr.text_frame.paragraphs[0].text = right_title
    right_hdr.text_frame.paragraphs[0].font.size = Pt(12)
    right_hdr.text_frame.paragraphs[0].font.bold = True
    right_hdr.text_frame.paragraphs[0].font.color.rgb = WHITE
    right_hdr.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Left code box
    left_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(2.0), Inches(4.5), Inches(5.0)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = LIGHT_GRAY
    left_box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    tf_left.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(left_code.split('\n')):
        if i == 0:
            p = tf_left.paragraphs[0]
        else:
            p = tf_left.add_paragraph()
        p.text = line
        p.font.size = Pt(9)
        p.font.name = 'Courier New'
        p.font.color.rgb = DARK_GRAY

    # Right code box
    right_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(2.0), Inches(4.5), Inches(5.0)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = LIGHT_GRAY
    right_box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    tf_right.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(right_code.split('\n')):
        if i == 0:
            p = tf_right.paragraphs[0]
        else:
            p = tf_right.add_paragraph()
        p.text = line
        p.font.size = Pt(9)
        p.font.name = 'Courier New'
        p.font.color.rgb = DARK_GRAY

    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide


def add_table_slide(prs, title, headers, rows, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.5),
        Inches(9), Inches(0.4 * num_rows + 0.3)
    )
    table = table_shape.table

    # Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE

    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = value
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GRAY

    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs)

    # Slide 2: Agenda
    add_content_slide(prs, "Agenda", [
        "1. Legacy COBOL System Overview",
        "2. Modernization Strategy",
        "3. Architecture Mapping",
        "4. Code Conversion (Side-by-Side)",
        "5. Testing Strategy & Results",
        "6. Deployment Pipeline",
        "7. Results & Benefits",
        "8. Next Steps",
    ], notes="This presentation walks through the complete process of converting a legacy COBOL accounting system to Node.js.")

    # Slide 3: Legacy System
    add_section_slide(prs, "Legacy COBOL System", "Understanding the existing architecture")

    # Slide 4: COBOL Architecture
    add_content_slide(prs, "Legacy COBOL Architecture", [
        "Three-module COBOL accounting system:",
        "  main.cob - User interface and menu loop (PERFORM UNTIL)",
        "  operations.cob - Business logic (credit, debit, view balance)",
        "  data.cob - In-memory data storage (READ/WRITE pattern)",
        "",
        "Key Characteristics:",
        "  Fixed-point arithmetic (PIC 9(6)V99 = 6 digits, 2 decimals)",
        "  Inter-program communication via CALL and LINKAGE SECTION",
        "  Sequential processing with EVALUATE (switch) statements",
        "  Compiled with GnuCOBOL to native executable",
    ], notes="The COBOL system uses a classic 3-tier architecture: UI layer, business logic, and data layer. Each is a separate program linked via CALL statements.")

    # Slide 5: COBOL Code Sample
    add_code_slide(prs,
        "COBOL Source Code Overview",
        "main.cob (UI Layer)",
        """IDENTIFICATION DIVISION.
PROGRAM-ID. MainProgram.

DATA DIVISION.
WORKING-STORAGE SECTION.
01 USER-CHOICE    PIC 9 VALUE 0.
01 CONTINUE-FLAG  PIC X(3) VALUE 'YES'.

PROCEDURE DIVISION.
MAIN-LOGIC.
  PERFORM UNTIL CONTINUE-FLAG = 'NO'
    DISPLAY "1. View Balance"
    DISPLAY "2. Credit Account"
    DISPLAY "3. Debit Account"
    DISPLAY "4. Exit"
    ACCEPT USER-CHOICE
    EVALUATE USER-CHOICE
      WHEN 1
        CALL 'Operations' USING 'TOTAL'
      WHEN 4
        MOVE 'NO' TO CONTINUE-FLAG
    END-EVALUATE
  END-PERFORM
  STOP RUN.""",
        "operations.cob (Logic)",
        """IDENTIFICATION DIVISION.
PROGRAM-ID. Operations.

DATA DIVISION.
WORKING-STORAGE SECTION.
01 AMOUNT        PIC 9(6)V99.
01 FINAL-BALANCE PIC 9(6)V99.

LINKAGE SECTION.
01 PASSED-OPERATION PIC X(6).

PROCEDURE DIVISION USING
    PASSED-OPERATION.

  IF OPERATION-TYPE = 'CREDIT'
    ACCEPT AMOUNT
    CALL 'DataProgram'
      USING 'READ', FINAL-BALANCE
    ADD AMOUNT TO FINAL-BALANCE
    CALL 'DataProgram'
      USING 'WRITE', FINAL-BALANCE
  END-IF
  GOBACK.""",
        notes="Left: The main program handles the menu loop using PERFORM UNTIL. Right: Operations handles business logic and communicates with the data layer via CALL statements."
    )

    # Slide 6: Modernization Strategy
    add_section_slide(prs, "Modernization Strategy", "Methodology and approach")

    # Slide 7: Approach
    add_content_slide(prs, "Migration Approach", [
        "Strategy: 1:1 Module-to-Class Mapping",
        "",
        "Principles:",
        "  Preserve exact business logic (no functional changes)",
        "  Maintain modular separation of concerns",
        "  Enable testability via dependency injection",
        "  Add REST API for cloud-native deployment",
        "",
        "Process:",
        "  1. Analyze COBOL data types and control flow",
        "  2. Map each PROGRAM-ID to a JavaScript class",
        "  3. Convert CALL statements to method invocations",
        "  4. Implement test suite from TESTPLAN.md",
        "  5. Containerize with Docker for deployment",
    ], notes="Our approach preserves the exact same business behavior while modernizing the technology stack. Each COBOL program becomes a JS class, and CALL statements become method calls.")

    # Slide 8: Architecture Mapping
    add_table_slide(prs, "Architecture Mapping (COBOL to Node.js)", [
        "COBOL Component", "Node.js Equivalent", "Pattern Change"
    ], [
        ["main.cob (MainProgram)", "src/main.js (AccountApp)", "PERFORM → while loop"],
        ["operations.cob", "src/operations.js (Operations)", "CALL → method call"],
        ["data.cob (DataProgram)", "src/data.js (DataStore)", "LINKAGE → constructor DI"],
        ["EVALUATE", "switch statement", "Same semantics"],
        ["PIC 9(6)V99", "number + toFixed(2)", "Floating point"],
        ["ACCEPT/DISPLAY", "readline / console", "Async I/O"],
        ["— (new)", "src/server.js", "HTTP REST API"],
    ], notes="Each COBOL module maps directly to a Node.js class. The key pattern change is from CALL/LINKAGE to constructor-based dependency injection.")

    # Slide 9: Code Conversion - Data Layer
    add_code_slide(prs,
        "Code Conversion: Data Layer",
        "data.cob (COBOL)",
        """IDENTIFICATION DIVISION.
PROGRAM-ID. DataProgram.

DATA DIVISION.
WORKING-STORAGE SECTION.
01 STORAGE-BALANCE PIC 9(6)V99
   VALUE 1000.00.

LINKAGE SECTION.
01 PASSED-OPERATION PIC X(6).
01 BALANCE          PIC 9(6)V99.

PROCEDURE DIVISION USING
    PASSED-OPERATION BALANCE.

  IF OPERATION-TYPE = 'READ'
    MOVE STORAGE-BALANCE TO BALANCE
  ELSE IF OPERATION-TYPE = 'WRITE'
    MOVE BALANCE TO STORAGE-BALANCE
  END-IF
  GOBACK.""",
        "data.js (Node.js)",
        """'use strict';

class DataStore {
  constructor(initialBalance = 1000.00) {
    this._balance = initialBalance;
  }

  read() {
    return this._balance;
  }

  write(newBalance) {
    this._balance = newBalance;
  }
}

module.exports = DataStore;




// COBOL LINKAGE SECTION replaced by
// constructor dependency injection
// PIC 9(6)V99 replaced by JS number
// IF 'READ'/'WRITE' replaced by
// named methods""",
        notes="The COBOL DataProgram uses string-based operations (READ/WRITE) via LINKAGE SECTION. In Node.js, this becomes a simple class with named methods. The STORAGE-BALANCE becomes a private instance variable."
    )

    # Slide 10: Code Conversion - Operations
    add_code_slide(prs,
        "Code Conversion: Business Logic",
        "operations.cob (COBOL)",
        """IF OPERATION-TYPE = 'DEBIT '
  DISPLAY "Enter debit amount: "
  ACCEPT AMOUNT
  CALL 'DataProgram'
    USING 'READ', FINAL-BALANCE
  IF FINAL-BALANCE >= AMOUNT
    SUBTRACT AMOUNT FROM
      FINAL-BALANCE
    CALL 'DataProgram'
      USING 'WRITE', FINAL-BALANCE
    DISPLAY "Amount debited."
      " New balance: "
      FINAL-BALANCE
  ELSE
    DISPLAY "Insufficient funds"
      " for this debit."
  END-IF
END-IF
GOBACK.""",
        "operations.js (Node.js)",
        """debit(amount) {
  const parsedAmount =
    parseFloat(amount);
  if (isNaN(parsedAmount) ||
      parsedAmount < 0) {
    return { success: false,
      balance: this.dataStore.read(),
      message: 'Invalid amount.' };
  }

  const currentBalance =
    this.dataStore.read();
  if (currentBalance >= parsedAmount) {
    const newBalance =
      currentBalance - parsedAmount;
    this.dataStore.write(newBalance);
    return { success: true,
      balance: newBalance,
      message: 'Amount debited. ' +
        'New balance: ' +
        newBalance.toFixed(2) };
  } else {
    return { success: false,
      balance: currentBalance,
      message: 'Insufficient funds' };
  }
}""",
        notes="The debit operation shows how COBOL's IF/ELSE maps to JavaScript. Note the addition of input validation (NaN check, negative check) that wasn't in the original COBOL. Return objects replace DISPLAY for better composability."
    )

    # Slide 11: Testing Strategy
    add_section_slide(prs, "Testing Strategy", "Comprehensive test coverage from TESTPLAN.md")

    # Slide 12: Test Results
    add_table_slide(prs, "Test Coverage Results", [
        "Test Case", "Description", "Status"
    ], [
        ["TC-1.1", "View Current Balance", "PASS"],
        ["TC-2.1", "Credit Account with Valid Amount", "PASS"],
        ["TC-2.2", "Credit Account with Zero Amount", "PASS"],
        ["TC-3.1", "Debit Account with Valid Amount", "PASS"],
        ["TC-3.2", "Debit with Amount > Balance (Insufficient Funds)", "PASS"],
        ["TC-3.3", "Debit Account with Zero Amount", "PASS"],
        ["TC-4.1", "Exit the Application", "PASS"],
    ], notes="All 7 test cases from the original TESTPLAN.md pass. Additional edge case tests (invalid inputs, negative amounts, NaN) bring total to 37 tests.")

    # Slide 13: Test Coverage Details
    add_content_slide(prs, "Test Coverage Metrics", [
        "37 total tests (unit + integration)",
        "",
        "Coverage Results:",
        "  Statements: 94.57%",
        "  Branches: 82.75%",
        "  Functions: 95.00%",
        "  Lines: 94.53%",
        "",
        "Test Structure:",
        "  tests/unit/data.test.js - DataStore (100% coverage)",
        "  tests/unit/operations.test.js - Business logic (100% coverage)",
        "  tests/unit/server.test.js - HTTP API endpoints",
        "  tests/integration/app.test.js - Full CLI flow simulation",
        "",
        "Framework: Jest with built-in coverage reporting",
    ], notes="We exceeded the 80% coverage threshold on all metrics. The data and operations modules have 100% coverage. The only uncovered lines are the main module's CLI entry point (process.exit handler).")

    # Slide 14: Deployment Pipeline
    add_section_slide(prs, "Deployment Pipeline", "Docker, CI/CD, and Cloud Deployment")

    # Slide 15: Deployment Architecture
    add_content_slide(prs, "Deployment Architecture", [
        "Multi-stage Docker Build:",
        "  Stage 1: Builder - Install dependencies (npm ci)",
        "  Stage 2: Production - Minimal Alpine image, non-root user",
        "",
        "Container Features:",
        "  Node.js 18 Alpine (minimal footprint)",
        "  Non-root user (security best practice)",
        "  Built-in health check endpoint (/health)",
        "  Resource limits: 128MB RAM, 0.5 CPU",
        "",
        "Deploy Script (deploy.sh) Commands:",
        "  full - Complete pipeline: lint, test, build, start",
        "  deploy-aws - Push to ECR, update ECS service",
        "  deploy-azure - Deploy to Azure App Service",
        "  health - Verify deployment via health endpoint",
    ], notes="The deployment uses a multi-stage Docker build for minimal image size. The deploy.sh script supports the full lifecycle from building to deploying to AWS or Azure.")

    # Slide 16: Deployment Commands
    add_code_slide(prs,
        "Deployment Script Usage",
        "Local Deployment",
        """# Full pipeline (recommended)
./deploy/deploy.sh full

# Step by step
./deploy/deploy.sh lint
./deploy/deploy.sh test
./deploy/deploy.sh build
./deploy/deploy.sh start

# Monitoring
./deploy/deploy.sh health
./deploy/deploy.sh logs

# Cleanup
./deploy/deploy.sh stop
./deploy/deploy.sh clean""",
        "Cloud Deployment",
        """# AWS ECS Deployment
export AWS_REGION=us-east-1
export ECR_REPO=node-accounting-app
export ECS_CLUSTER=app-cluster
export ECS_SERVICE=app-service
./deploy/deploy.sh deploy-aws

# Azure App Service
export AZURE_RESOURCE_GROUP=app-rg
export AZURE_WEBAPP_NAME=app
./deploy/deploy.sh deploy-azure

# Health Check
curl http://localhost:3000/health
# {"status":"healthy",
#  "service":"node-accounting-app",
#  "version":"1.0.0"}""",
        notes="The deploy script provides both local Docker deployment and cloud deployment to AWS ECS or Azure App Service. A single command handles the full pipeline."
    )

    # Slide 17: Results & Benefits
    add_section_slide(prs, "Results & Benefits", "Modernization outcomes")

    # Slide 18: Benefits
    add_table_slide(prs, "Modernization Benefits", [
        "Category", "Before (COBOL)", "After (Node.js)"
    ], [
        ["Language", "COBOL (1959)", "JavaScript/Node.js (2009+)"],
        ["Developers", "Scarce, expensive", "Abundant, accessible"],
        ["Testing", "Manual only", "Automated (37 tests, 94% coverage)"],
        ["Deployment", "Compile & copy", "Docker + CI/CD pipeline"],
        ["Cloud Ready", "No", "AWS ECS, Azure App Service"],
        ["API Access", "Terminal only", "REST API + CLI"],
        ["Monitoring", "None", "Health checks, structured logs"],
        ["Extensibility", "Difficult", "npm ecosystem, modular classes"],
    ], notes="The modernization delivers immediate benefits in testability, deployability, and developer accessibility while preserving 100% of the original business logic.")

    # Slide 19: Next Steps
    add_content_slide(prs, "Next Steps & Recommendations", [
        "Immediate:",
        "  Set up CI/CD pipeline (GitHub Actions / Jenkins)",
        "  Deploy to staging environment",
        "  Performance benchmarking",
        "",
        "Short-term:",
        "  Add persistent database (PostgreSQL/MongoDB)",
        "  Implement authentication & authorization",
        "  Add transaction history and audit logging",
        "",
        "Long-term:",
        "  Multi-account support",
        "  WebSocket real-time updates",
        "  Microservice decomposition (if needed)",
        "  API Gateway integration",
    ], notes="The modular architecture makes these enhancements straightforward. The DataStore can be swapped for a database by implementing the same read/write interface.")

    # Slide 20: Summary
    add_content_slide(prs, "Summary", [
        "Successfully converted 3-module COBOL system to Node.js",
        "",
        "Deliverables:",
        "  node-app/ - Complete Node.js application (4 modules)",
        "  tests/ - 37 automated tests with 94%+ coverage",
        "  docs/ - Architecture and deployment documentation",
        "  deploy/ - Docker + cloud deployment scripts",
        "",
        "Key Metrics:",
        "  100% feature parity with original COBOL system",
        "  All TESTPLAN.md test cases passing",
        "  Zero runtime dependencies (Node.js built-ins only)",
        "  Production-ready Docker image with health checks",
        "",
        "Questions?",
    ], notes="The modernization is complete and production-ready. All original functionality is preserved with significant improvements in testability, deployability, and maintainability.")

    # Save the presentation
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, 'COBOL_to_NodeJS_Migration.pptx')
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == '__main__':
    main()
