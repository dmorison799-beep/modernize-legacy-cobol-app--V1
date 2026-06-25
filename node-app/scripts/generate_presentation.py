#!/usr/bin/env python3
"""Generate a PowerPoint presentation summarizing the COBOL to Node.js migration."""

import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)


def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "Welcome to the COBOL to Node.js migration overview. "
        "This presentation covers the modernization of a legacy "
        "accounting system from COBOL to Node.js."
    )
    return slide


def add_content_slide(prs, title, bullets, notes_text=""):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for i, bullet in enumerate(bullets):
        if i == 0:
            tf.text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

    if notes_text:
        notes = slide.notes_slide
        notes.notes_text_frame.text = notes_text

    return slide


def create_presentation():
    prs = Presentation()

    # Slide 1: Title
    add_title_slide(
        prs,
        "COBOL to Node.js Migration",
        "Account Management System Modernization"
    )

    # Slide 2: Project Overview
    add_content_slide(
        prs,
        "Project Overview",
        [
            "Legacy COBOL accounting system converted to Node.js",
            "Maintains identical business logic and user experience",
            "Interactive CLI for balance viewing, credits, and debits",
            "In-memory data persistence with default $1,000 balance",
            "Modular architecture matching original COBOL structure"
        ],
        "The original COBOL application was a terminal-based account management "
        "system. The Node.js version preserves the same functionality while "
        "enabling modern development practices."
    )

    # Slide 3: Architecture Mapping
    add_content_slide(
        prs,
        "Architecture: COBOL to Node.js Mapping",
        [
            "main.cob -> src/main.js (CLI entry point with readline)",
            "operations.cob -> src/operations.js (credit, debit, getBalance)",
            "data.cob -> src/data.js (in-memory data persistence)",
            "CALL/GOBACK -> module.exports / require()",
            "PIC 9(6)V99 -> JavaScript Number with toFixed(2)"
        ],
        "Each COBOL program maps directly to a Node.js module. "
        "The CALL/GOBACK pattern is replaced with CommonJS module imports. "
        "COBOL's fixed-point decimal is handled with JavaScript Number "
        "and explicit 2-decimal rounding."
    )

    # Slide 4: Key Features Preserved
    add_content_slide(
        prs,
        "Key Features Preserved",
        [
            "View current account balance",
            "Credit account with specified amount",
            "Debit account with insufficient funds validation",
            "Interactive menu-driven interface",
            "Default starting balance of $1,000.00"
        ],
        "All core business logic from the COBOL system is preserved. "
        "The insufficient funds check ensures debits cannot exceed "
        "the current balance, matching the original COBOL behavior."
    )

    # Slide 5: Testing Strategy
    add_content_slide(
        prs,
        "Testing Strategy",
        [
            "Jest testing framework with 7 test cases from TESTPLAN.md",
            "Unit tests for data persistence layer (data.test.js)",
            "Unit tests for business logic (operations.test.js)",
            "Integration tests for end-to-end flows (integration.test.js)",
            "Test cases: TC-1.1, TC-2.1, TC-2.2, TC-3.1, TC-3.2, TC-3.3, TC-4.1"
        ],
        "The test suite covers all scenarios from the original TESTPLAN.md, "
        "including edge cases like zero amounts and insufficient funds. "
        "Tests are organized into unit and integration layers."
    )

    # Slide 6: Deployment
    add_content_slide(
        prs,
        "Deployment",
        [
            "Bash deployment script: deploy.sh [environment]",
            "Supports development, staging, and production environments",
            "Automated validation: dependency install, lint, and test",
            "Environment-specific configuration via NODE_ENV",
            "Zero-downtime deployment pipeline ready"
        ],
        "The deployment script automates the full deployment pipeline "
        "including dependency installation, linting, testing, and "
        "environment-specific configuration."
    )

    # Slide 7: Benefits of Modernization
    add_content_slide(
        prs,
        "Benefits of Modernization",
        [
            "Modern language with large ecosystem (npm)",
            "Easier to find and train developers",
            "Comprehensive automated testing with Jest",
            "ESLint for code quality enforcement",
            "Ready for future enhancements (REST API, database, UI)"
        ],
        "Moving from COBOL to Node.js opens up significant opportunities "
        "for future development. The modular architecture makes it easy "
        "to add a REST API, connect to a database, or build a web frontend."
    )

    # Slide 8: Next Steps
    add_content_slide(
        prs,
        "Next Steps",
        [
            "Validate business logic with stakeholders",
            "Add persistent storage (database integration)",
            "Build REST API layer for web/mobile access",
            "Implement authentication and authorization",
            "Set up CI/CD pipeline for automated deployments"
        ],
        "These next steps will further modernize the application "
        "and prepare it for enterprise-grade deployment."
    )

    return prs


def main():
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output_dir = os.path.join(script_dir, "..")
    output_path = os.path.join(output_dir, "COBOL_to_NodeJS_Migration.pptx")

    print("Generating COBOL to Node.js Migration presentation...")
    prs = create_presentation()
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")


if __name__ == "__main__":
    main()
